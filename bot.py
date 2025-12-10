# bot.py - Lotinify Telegram Boti (YAKUNIY BARQAROR VA INLINE VERSIYA)

import asyncio
import logging
import io
import re
from typing import Tuple, Union, Any

# Aiogram kutubxonalari
from aiogram import Bot, Dispatcher, types, F, Router
from aiogram.enums import ParseMode
from aiogram.filters import Command
from aiogram.types import ReplyKeyboardMarkup, KeyboardButton, BufferedInputFile
from aiogram.types import InlineQueryResultArticle, InputTextMessageContent, InlineQuery
from aiogram.utils.markdown import html_decoration
from aiogram.fsm.context import FSMContext
from aiogram.fsm.state import State, StatesGroup

# --- AI API himoyalangan import (google-genai) ---
genai = None
gemini_types = None
class APIError(Exception): pass 

try:
    from google import genai
    from google.genai import types as gemini_types
    from google.genai import errors as gemini_errors 
    APIError = gemini_errors.APIError 
except ImportError:
    logging.warning("Gemini API kutubxonasi 'google-genai' topilmadi. Imlo tekshirish ISHLAMAYDI.")
except Exception as e:
    logging.error(f"Gemini Clientni yuklashda kutilmagan xato: {e}")

# --- Office fayl kutubxonalari himoyalangan import (python-docx, openpyxl, python-pptx) ---
try:
    from docx import Document
    from openpyxl import load_workbook
    from pptx import Presentation
except ImportError:
    logging.warning("Office fayl kutubxonalari topilmadi. Fayl konvertatsiyasi ISHLAMAYDI.")
    Document = None
    load_workbook = None
    Presentation = None

# ---------------- LOG va KONFIGURATSIYA ----------------
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# !!! YANGILANGAN BOT TOKENDAN FOYDALANILMOQDA !!!
TOKEN = "8420487214:AAFgPefTZNiF843hOZjYb_-3J_6V6SuYzmY" 
ADMIN_ID = 1455902088 
GEMINI_API_KEY = "AIzaSyBh74Q6js175gbAZ3muhphrmjfO-KOB8qU" 

MAX_FILE_SIZE_MB = 50
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024
GEMINI_MODEL = "gemini-2.5-flash"

bot = Bot(token=TOKEN)
dp = Dispatcher()
router = Router()

# Gemini mijozini yaratish
gemini_client = None
if GEMINI_API_KEY and genai:
    try:
        gemini_client = genai.Client(api_key=GEMINI_API_KEY)
    except Exception as e:
        logger.error(f"Gemini Clientni ishga tushirishda xato: {e}")
        gemini_client = None


# ======================
# FSM States & Menus
# ======================
class TranslitState(StatesGroup):
    waiting_for_text = State()
    waiting_for_spellcheck = State()
    waiting_for_menu = State()
    waiting_for_docx = State()
    waiting_for_xlsx = State()
    waiting_for_pptx = State()

MAIN_MENU = ReplyKeyboardMarkup(keyboard=[
    [KeyboardButton(text="üîÑ Lotin ‚Üî Kirill (Avtomatik)"), KeyboardButton(text="‚úèÔ∏è Imlo Tekshirish (Gemini)")],
    [KeyboardButton(text="üìÅ Fayl Konvertatsiyasi")],
], resize_keyboard=True)

FILE_MENU = ReplyKeyboardMarkup(keyboard=[
    [KeyboardButton(text="üìÑ Word (.docx)"), KeyboardButton(text="üìä Excel (.xlsx)")],
    [KeyboardButton(text="üé¨ PowerPoint (.pptx)")],
    [KeyboardButton(text="‚óÄÔ∏è Orqaga")]
], resize_keyboard=True)

# ==========================================================
# Lotin ‚Üî Kirill Mantiqi
# ==========================================================

def normalize_apostrophes(text: str) -> str:
    bad_apos = ["‚Äú", "‚Äù", "‚Äò", "‚Äô", " ª", " º", "`", "¬¥", " π", " æ", "Àà", " Ω"]
    for b in bad_apos:
        text = text.replace(b, "'")
    return text

# --- 1. KIRILLDAN LOTINGA ---
CYR_TO_LAT_MAP = {
    '—û': "o ª", '“ì': "g ª", '“õ': "q", '“≥': "h", '—à': "sh", '—á': "ch", '—Ü': "ts", '—ä': "'", '—å': '', '—ç': 'e',
    '–é': "O ª", '“í': "G ª", '“ö': "Q", '“≤': "H", '–®': "Sh", '–ß': "Ch", '–¶': "Ts", '–™': "'", '–¨': '', '–≠': 'E',
    '–∞': 'a', '–±': 'b', '–≤': 'v', '–≥': 'g', '–¥': 'd', '–µ': 'e', '–∂': 'j', '–∑': 'z', '–∏': 'i', '–π': 'y', '–∫': 'k', '–ª': 'l', '–º': 'm', '–Ω': 'n', 
    '–æ': 'o', '–ø': 'p', '—Ä': 'r', '—Å': 's', '—Ç': 't', '—É': 'u', '—Ñ': 'f', '—Ö': 'x', '—è': 'ya', '—é': 'yu', '—ë': 'yo',
    '–ê': 'A', '–ë': 'B', '–í': 'V', '–ì': 'G', '–î': 'D', '–ï': 'Ye', '–ñ': 'J', '–ó': 'Z', '–ò': 'I', '–ô': 'Y', '–ö': 'K', '–õ': 'L', '–ú': 'M', '–ù': 'N', 
    '–û': 'O', '–ü': 'P', '–†': 'R', '–°': 'S', '–¢': 'T', '–£': 'U', '–§': 'F', '–•': 'X', '–Ø': 'Ya', '–Æ': 'Yu', '–Å': 'Yo',
}

def uz_to_latin(text: str) -> str:
    text = normalize_apostrophes(text)
    
    sorted_cyr_keys = sorted(CYR_TO_LAT_MAP.keys(), key=len, reverse=True)
    
    for cyr in sorted_cyr_keys:
        text = text.replace(cyr, CYR_TO_LAT_MAP[cyr])

    def fix_initial_ye(match):
        return 'Ye' if match.group(0) == '–ï' else match.group(0)
    
    return re.sub(r'(^|\s|[^a-zA-Z])([–ï])', fix_initial_ye, text)

# --- 2. LOTINDAN KIRILLGA ---
def uz_to_cyrillic(text: str) -> str:
    text = normalize_apostrophes(text)
    result = []
    i = 0
    
    CYR_MAP = {
        'a': '–∞', 'b': '–±', 'v': '–≤', 'g': '–≥', 'd': '–¥', 'z': '–∑', 'i': '–∏', 'y': '–π', 'k': '–∫', 'l': '–ª',
        'm': '–º', 'n': '–Ω', 'o': '–æ', 'p': '–ø', 'r': '—Ä', 's': '—Å', 't': '—Ç', 'u': '—É', 'f': '—Ñ',
        'x': '—Ö', 'h': '“≥', 'q': '“õ', 'j': '–∂', "'": '—ä'
    }

    while i < len(text):
        char = text[i]
        
        is_upper = char.isupper()
        current_part = text[i:].lower() 
        current_char_lower = char.lower()

        found = False
        
        compound_map = {
            'o ª': '—û', 'g ª': '“ì', 'sh': '—à', 'ch': '—á', 'ng': '–Ω–≥',
            "o'": '—û', "g'": '“ì',
        }

        for lat, cyr in sorted(compound_map.items(), key=lambda item: len(item[0]), reverse=True):
            if current_part.startswith(lat):
                if is_upper and len(lat) > 1 and lat.islower():
                     result.append(cyr.upper())
                elif char.isupper() and len(lat) == 2 and text[i+1].islower():
                     result.append(cyr.capitalize())
                else:
                     result.append(cyr.upper() if is_upper else cyr)

                i += len(lat)
                found = True
                break
        
        if found:
            continue
            
        if current_part.startswith('ts'):
             is_at_start_of_word = (i == 0 or not text[i-1].isalpha())
             
             if is_at_start_of_word:
                 result.append('–¶' if is_upper else '—Ü')
                 i += 2
                 continue
             else:
                 pass

        if current_char_lower == 'y' and i + 1 < len(text) and text[i+1].lower() in ('a', 'o', 'u', 'e'):
            next_char = text[i+1].lower()
            
            cyr_map_y_vowel = {'a': '—è', 'o': '—ë', 'u': '—é', 'e': '–µ'}
            cyr = cyr_map_y_vowel[next_char]
            
            result.append(cyr.upper() if is_upper else cyr)
            i += 2
            continue
            
        if current_char_lower == 'e':
            is_start_of_word = (i == 0 or not text[i-1].isalpha())
            
            if is_start_of_word:
                result.append('–≠' if is_upper else '—ç')
            else:
                result.append('–ï' if is_upper else '–µ') 
            i += 1
            continue
            
        if current_char_lower in CYR_MAP:
            cyr = CYR_MAP[current_char_lower]
            result.append(cyr.upper() if is_upper else cyr)
        else:
            result.append(char)
        
        i += 1
        
    final_text = "".join(result)
    
    final_text = final_text.replace("–π–∞", "—è").replace("–π–æ", "—ë").replace("–π—É", "—é").replace("–π–µ", "–µ")
    final_text = final_text.replace("–π–µ—Ä", "–µ—Ä") 
    final_text = final_text.replace('—ë—ä', '–π—û') 
    
    return final_text

# --- Skriptni aniqlash funksiyasi ---
def detect_script(text: str) -> str:
    cyr_count = 0
    lat_count = 0
    t = text.lower()
    
    cyr_count += t.count('—û') * 10 + t.count('“ì') * 10 + t.count('“õ') * 5 + t.count('“≥') * 5
    lat_count += t.count('o ª') * 10 + t.count('g ª') * 10
    
    lat_count += sum(t.count(c) for c in ["sh", "ch", "ts", "ya", "yu", "yo"]) * 3
    cyr_count += sum(t.count(c) for c in ["—à", "—á", "—Ü", "—è", "—é", "—ë"]) * 3
    
    for char in t:
        if '–∞' <= char <= '—è':
            cyr_count += 1
        elif 'a' <= char <= 'z':
            lat_count += 1
            
    if cyr_count > lat_count * 2:
        return 'cyrillic'
    elif lat_count > cyr_count * 2:
        return 'latin'
    else:
        return 'unknown'
        
# ==========================================================
# Handlers va Fayl Mantiqi
# ==========================================================

# --- GEMINI IMMLO TEKSHIRISH FUNKSIYASI ---
async def gemini_process_text(text: str, task_type: str) -> str:
    if gemini_client is None:
        return "‚ùå Gemini API sozlanmagan. Imloni tekshirib bo'lmadi."
    
    detected_script = detect_script(text)
    
    if detected_script == 'cyrillic':
        target_script = "KIRILL O'ZBEK ALIFBOSIDAN"
    else:
        target_script = "LOTIN O'ZBEK ALIFBOSIDAN"
        
    if task_type == 'spellcheck':
        system_instruction = (
            "Siz professional O'zbek tili imlo va grammatika tekshiruvchisiz. "
            "Berilgan matnni tahlil qiling. Barcha imlo va tinish belgisi xatolarini to'g'rilang. "
            f"JAVOBINGIZDA MUTLAQO VA FAQA–¢ {target_script} FOYDALANING. "
            "Faqat va faqat to'g'rilangan matnni qaytaring. Hech qanday qo'shimcha izoh yozmang."
        )
    else:
        return "Noma'lum vazifa turi."
        
    try:
        config = gemini_types.GenerateContentConfig(
            system_instruction=system_instruction,
            temperature=0.0,
        )
        
        # <<< ASOSIY TUZATISH: generate_content sinxron funksiyasini asyncio.to_thread orqali chaqirish >>>
        response = await asyncio.to_thread(
            gemini_client.models.generate_content, 
            model=GEMINI_MODEL,
            contents=[text],
            config=config
        )
        
        result_text = response.text.strip()
        
        # Qaytarilgan matnni foydalanuvchi kiritgan alifboga mosligini tekshirish (fallback)
        if detected_script == 'latin' and detect_script(result_text) == 'cyrillic':
             logger.warning("Gemini Lotincha so'ralganda Kirillcha javob qaytardi, lotinchaga o'tkazilmoqda.")
             result_text = uz_to_latin(result_text)
        elif detected_script == 'cyrillic' and detect_script(result_text) == 'latin':
             logger.warning("Gemini Kirillcha so'ralganda Lotincha javob qaytardi, kirillchaga o'tkazilmoqda.")
             result_text = uz_to_cyrillic(result_text)
             
        return result_text
        
    except APIError as e: 
        logger.error(f"Gemini API xatosi (Ulanish): {e}")
        return f"‚ùå Gemini API xatosi. Iltimos, server ulanishini yoki API kalitini tekshiring. Xato: {html_decoration.quote(str(e))}"
        
    except Exception as e:
        logger.exception(f"Gemini bilan bog'lanishda kutilmagan ichki xato yuz berdi: {e}")
        return f"‚ùå Gemini bilan bog'lanishda kutilmagan ichki xato yuz berdi. (Tafsilotlar logda: {html_decoration.quote(str(e))})"
        
# --- Fayl konvertatsiya mantig'i ---
async def convert_office_file_from_bytes(file_bytes: io.BytesIO, filename: str) -> Union[Tuple[io.BytesIO, str], str]:
    if not (Document and load_workbook and Presentation):
         return "‚ùå Office fayllar bilan ishlash kutubxonalari o'rnatilmagan."
         
    ext = filename.lower().split('.')[-1]
    file_bytes.seek(0)
    
    if 'kirill' in filename.lower():
          new_filename = filename.replace('kirill', 'lotin')
    elif 'lotin' in filename.lower():
          new_filename = filename.replace('lotin', 'kirill')
    else:
          new_filename = f"converted_{filename}"

    try:
        if ext == 'docx':
            doc = Document(file_bytes)
            paragraphs = list(doc.paragraphs)
            for table in doc.tables:
                 for row in table.rows:
                     for cell in row.cells:
                         paragraphs.extend(cell.paragraphs)
            
            sample = "\n".join(p.text for p in paragraphs if p.text)
            script = detect_script(sample)
            
            if script == 'unknown': return "Word: Alifbo turi aniqlanmadi."
            converter_func = uz_to_latin if script == 'cyrillic' else uz_to_cyrillic

            for para in paragraphs:
                if para.text and para.text.strip():
                    new_text = converter_func(para.text)
                    if para.runs:
                        para.runs[0].text = new_text
                        for run in para.runs[1:]:
                            run.text = "" 
                    else:
                        para.text = new_text

            out = io.BytesIO()
            doc.save(out)
            out.seek(0)
            return out, new_filename
        
        elif ext == 'xlsx':
            wb = load_workbook(file_bytes)
            script = None
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell_value in row:
                        if isinstance(cell_value, str) and cell_value.strip():
                            script = detect_script(cell_value)
                            if script in ['cyrillic', 'latin']: break
                    if script in ['cyrillic', 'latin']: break
                if script in ['cyrillic', 'latin']: break

            if not script or script == 'unknown': return "Excel: Alifbo turi aniqlanmadi."
            converter_func = uz_to_latin if script == 'cyrillic' else uz_to_cyrillic

            for sheet in wb.worksheets:
                for row in sheet.iter_rows():
                    for cell in row:
                        if isinstance(cell.value, str) and cell.value.strip():
                            cell.value = converter_func(cell.value)
            
            out = io.BytesIO()
            wb.save(out)
            out.seek(0)
            return out, new_filename
            
        elif ext == 'pptx':
            prs = Presentation(file_bytes)
            text_sample = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame: continue
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text and paragraph.text.strip():
                            text_sample += paragraph.text + " "
                            
            script = detect_script(text_sample) if text_sample else 'unknown'
            
            if script == 'unknown': return "PowerPoint: Alifbo turi aniqlanmadi."

            converter_func = uz_to_latin if script == 'cyrillic' else uz_to_cyrillic

            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame: continue
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text and paragraph.text.strip():
                            paragraph.text = converter_func(paragraph.text)

            out = io.BytesIO()
            prs.save(out)
            out.seek(0)
            return out, new_filename
        
        else:
             return f"Bot faqat **.docx, .xlsx, .pptx** formatlarini qo'llab-quvvatlaydi."
    
    except Exception as e:
        logger.exception(f"Office conversion failed for .{ext}")
        return f"Ichki texnik xatolik: Fayl buzilgan bo'lishi mumkin."


async def process_file_logic(msg: types.Message, state: FSMContext, expected_ext: str):
    document = msg.document
    
    if document.file_size and document.file_size > MAX_FILE_SIZE_BYTES:
        await msg.answer(f"‚ùå Fayl hajmi limiti ({MAX_FILE_SIZE_MB} MB) dan oshib ketdi.", reply_markup=FILE_MENU)
        return

    processing = await msg.answer("üìÇ Fayl qabul qilindi, konvertatsiya jarayoni boshlandi...")

    try:
        file_bytes = io.BytesIO()
        await bot.download(file=document.file_id, destination=file_bytes)
        file_bytes.seek(0)
        
        result = await convert_office_file_from_bytes(file_bytes, document.file_name)

        await bot.delete_message(msg.chat.id, processing.message_id)

        if isinstance(result, str):
            await msg.answer(f"‚ùå Xatolik: <b>{html_decoration.quote(result)}</b>",
                             reply_markup=MAIN_MENU, parse_mode=ParseMode.HTML)
        else:
            out_io, new_name = result
            out_io.seek(0)
            
            escaped_new_name = html_decoration.quote(new_name) 

            await msg.answer_document(
                document=BufferedInputFile(out_io.read(), filename=new_name),
                caption=f"‚úÖ Konvertatsiya yakunlandi: <b>{escaped_new_name}</b>",
                reply_markup=MAIN_MENU,
                parse_mode=ParseMode.HTML 
            )

    except Exception as e:
        logger.exception(f"process_file_logic failed during download/send: {e}")
        await msg.answer(f"‚ùå Yuklab olish yoki yuborishda umumiy xatolik. Iltimos, yana urinib ko'ring.",
                         reply_markup=MAIN_MENU, parse_mode=ParseMode.MARKDOWN)
    finally:
        await state.clear()


# ==========================================================
# Handlers: Oddiy Rejim (Foydalanuvchi Bilan Chat)
# ==========================================================

@router.message(Command("start"))
async def start_handler(msg: types.Message, state: FSMContext):
    await state.clear()
    await msg.answer("Assalomu alaykum! Lotin ‚Üî Kirill konvertoriga xush kelibsiz.\n\n**Asosiy menyu:**", reply_markup=MAIN_MENU, parse_mode=ParseMode.MARKDOWN)

@router.message(F.text == "üîÑ Lotin ‚Üî Kirill (Avtomatik)")
async def auto_translit_entry(msg: types.Message, state: FSMContext):
    await state.set_state(TranslitState.waiting_for_text)
    await msg.answer("Iltimos, matn yuboring. Kiritilgan alifbo avtomatik aniqlanadi.",
                     reply_markup=ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text="‚óÄÔ∏è Orqaga")]], resize_keyboard=True), parse_mode=ParseMode.MARKDOWN)


@router.message(F.text == "‚úèÔ∏è Imlo Tekshirish (Gemini)")
async def spellcheck_entry(msg: types.Message, state: FSMContext):
    if gemini_client is None:
        await msg.answer("‚ùå Gemini API sozlanmaganligi sababli Imlo tekshirish funksiyasi ishlamaydi.", reply_markup=MAIN_MENU)
        return
        
    await state.set_state(TranslitState.waiting_for_spellcheck)
    await msg.answer("Iltimos, tekshirilishi kerak bo'lgan matnni yuboring. **Gemini** imlo va grammatik xatolarni to'g'rilaydi.",
                     reply_markup=ReplyKeyboardMarkup(keyboard=[[KeyboardButton(text="‚óÄÔ∏è Orqaga")]], resize_keyboard=True), parse_mode=ParseMode.MARKDOWN)


@router.message(F.text == "üìÅ Fayl Konvertatsiyasi")
async def file_translit_entry(msg: types.Message, state: FSMContext):
    if not (Document and load_workbook and Presentation):
         await msg.answer("‚ùå Fayl kutubxonalari o'rnatilmagan. Fayl konvertatsiyasi ISHLAMAYDI.", reply_markup=MAIN_MENU)
         return
         
    await state.set_state(TranslitState.waiting_for_menu)
    await msg.answer(f"Konvertatsiya turini tanlang. (Maksimal fayl hajmi: {MAX_FILE_SIZE_MB} MB)", reply_markup=FILE_MENU)

@router.message(F.text == "üìÑ Word (.docx)")
async def select_docx(msg: types.Message, state: FSMContext):
    await state.set_state(TranslitState.waiting_for_docx)
    await msg.answer(f"Iltimos, **.docx** faylini yuboring.", parse_mode=ParseMode.MARKDOWN, reply_markup=FILE_MENU)

@router.message(F.text == "üìä Excel (.xlsx)")
async def select_xlsx(msg: types.Message, state: FSMContext):
    await state.set_state(TranslitState.waiting_for_xlsx)
    await msg.answer(f"Iltimos, **.xlsx** faylini yuboring.", parse_mode=ParseMode.MARKDOWN, reply_markup=FILE_MENU)

@router.message(F.text == "üé¨ PowerPoint (.pptx)")
async def select_pptx(msg: types.Message, state: FSMContext):
    await state.set_state(TranslitState.waiting_for_pptx)
    await msg.answer(f"Iltimos, **.pptx** faylini yuboring.", parse_mode=ParseMode.MARKDOWN, reply_markup=FILE_MENU)

@router.message(F.text == "‚óÄÔ∏è Orqaga")
async def exit_mode_handler(msg: types.Message, state: FSMContext):
    current_state = await state.get_state()
    if current_state in [TranslitState.waiting_for_menu, TranslitState.waiting_for_text, TranslitState.waiting_for_spellcheck]:
        await state.clear()
        await msg.answer("Asosiy menyuga qaytildi:", reply_markup=MAIN_MENU)
    elif current_state in [TranslitState.waiting_for_docx, TranslitState.waiting_for_xlsx, TranslitState.waiting_for_pptx]:
        await state.set_state(TranslitState.waiting_for_menu)
        await msg.answer(f"Fayl turini tanlang.", reply_markup=FILE_MENU)
    else:
        await state.clear()
        await msg.answer("Asosiy menyu:", reply_markup=MAIN_MENU)

@router.message(TranslitState.waiting_for_docx, F.document)
async def process_docx_file(msg: types.Message, state: FSMContext):
    await process_file_logic(msg, state, 'docx')

@router.message(TranslitState.waiting_for_xlsx, F.document)
async def process_xlsx_file(msg: types.Message, state: FSMContext):
    await process_file_logic(msg, state, 'xlsx')

@router.message(TranslitState.waiting_for_pptx, F.document)
async def process_pptx_file(msg: types.Message, state: FSMContext):
    await process_file_logic(msg, state, 'pptx')

@router.message(TranslitState.waiting_for_text, F.text)
async def auto_convert_text_rule(msg: types.Message):
    input_text = msg.text
    script = detect_script(input_text)

    if script == 'cyrillic':
        result = uz_to_latin(input_text)
    elif script == 'latin':
        result = uz_to_cyrillic(input_text)
    else:
        await msg.answer("Alifbo aniqlanmadi. Lotin yoki Kirill yozuvidagi matn yuboring.")
        return

    await msg.answer(result, parse_mode=ParseMode.MARKDOWN)


@router.message(TranslitState.waiting_for_spellcheck, F.text)
async def process_spellcheck_gemini(msg: types.Message):
    input_text = msg.text

    # Kutish xabari
    checking_msg = await msg.answer("‚è≥ Matn **Gemini** orqali imlo va grammatik xatolar uchun tekshirilmoqda... \n\n_Bu bir necha soniya olishi mumkin._")
    
    result = await gemini_process_text(input_text, 'spellcheck') 
    
    # Xabarni natija bilan almashtirish
    await bot.edit_message_text(
        chat_id=msg.chat.id, 
        message_id=checking_msg.message_id, 
        text=result,
        parse_mode=ParseMode.MARKDOWN
    )

@router.message(F.text)
async def handle_unhandled_text(msg: types.Message, state: FSMContext):
    current_state = await state.get_state()
    
    if current_state in [TranslitState.waiting_for_docx, TranslitState.waiting_for_xlsx, TranslitState.waiting_for_pptx]:
        await msg.answer(f"Iltimos, fayl yuboring yoki **‚óÄÔ∏è Orqaga** tugmasini bosing.", reply_markup=FILE_MENU, parse_mode=ParseMode.MARKDOWN)
        return
    
    if current_state == TranslitState.waiting_for_menu:
        await msg.answer("Iltimos, menyudan fayl turini tanlang.", reply_markup=FILE_MENU)
        return
        
    await msg.answer("Iltimos, menyu tugmalaridan foydalaning.", reply_markup=MAIN_MENU)

@router.message(F.document)
async def handle_unexpected_file(msg: types.Message, state: FSMContext):
    current_state = await state.get_state()
    
    if current_state in [TranslitState.waiting_for_docx, TranslitState.waiting_for_xlsx, TranslitState.waiting_for_pptx]:
        expected_ext = current_state.state.split('_')[-1]
        await msg.answer(f"‚ùå Kutilmagan fayl turi. Iltimos, **.{expected_ext}** faylini yuboring.", reply_markup=FILE_MENU, parse_mode=ParseMode.MARKDOWN)
    else:
        await state.clear()
        await msg.answer("Menyu tugmalaridan foydalaning. Fayl yuborish uchun **üìÅ Fayl Konvertatsiyasi** tugmasini bosing.", reply_markup=MAIN_MENU)


# ==========================================================
# Handlers: INLINE Rejim
# ==========================================================

@router.inline_query()
async def inline_translit_handler(inline_query: InlineQuery):
    query_text = inline_query.query.strip()
    results = []

    if not query_text:
        results.append(
            InlineQueryResultArticle(
                id="help",
                title="Qo'llanma: @botname [lotin | kirill] [matn]",
                input_message_content=InputTextMessageContent(
                    message_text="Lotin ‚Üî Kirill konvertatsiyasi uchun shunday yozing:\n`@lotinifybot lotin matn`\nyoki\n`@lotinifybot kirill matn`",
                    parse_mode=ParseMode.MARKDOWN
                )
            )
        )
    else:
        parts = query_text.split(maxsplit=1)
        command = parts[0].lower() if parts else ""
        text_to_convert = parts[1] if len(parts) > 1 else ""

        if not text_to_convert:
            results.append(
                 InlineQueryResultArticle(
                    id="enter_text",
                    title="Matnni kiriting...",
                    input_message_content=InputTextMessageContent(
                        message_text="Matn kiritilmadi. Lotin yoki Kirill matnini yozing."
                    )
                )
            )
        elif command == "lotin":
            converted_text = uz_to_latin(text_to_convert)
            results.append(
                InlineQueryResultArticle(
                    id="to_latin",
                    title="‚úÖ Lotinchaga o'tkazish",
                    description=converted_text,
                    input_message_content=InputTextMessageContent(
                        message_text=converted_text
                    )
                )
            )
        elif command == "kirill":
            converted_text = uz_to_cyrillic(text_to_convert)
            results.append(
                InlineQueryResultArticle(
                    id="to_cyrillic",
                    title="‚úÖ Kirillchaga o'tkazish",
                    description=converted_text,
                    input_message_content=InputTextMessageContent(
                        message_text=converted_text
                    )
                )
            )
        else:
            # Agar komanda kiritilmagan bo'lsa, avtomatik aniqlashga harakat qilamiz
            script = detect_script(query_text)
            
            if script == 'cyrillic':
                converted_latin = uz_to_latin(query_text)
                results.append(
                    InlineQueryResultArticle(
                        id="auto_to_latin",
                        title="‚û°Ô∏è Lotinchaga o'tkazish",
                        description=converted_latin,
                        input_message_content=InputTextMessageContent(
                            message_text=converted_latin
                        )
                    )
                )
            elif script == 'latin':
                converted_cyr = uz_to_cyrillic(query_text)
                results.append(
                    InlineQueryResultArticle(
                        id="auto_to_cyrillic",
                        title="‚¨ÖÔ∏è Kirillchaga o'tkazish",
                        description=converted_cyr,
                        input_message_content=InputTextMessageContent(
                            message_text=converted_cyr
                        )
                    )
                )
            else:
                 results.append(
                    InlineQueryResultArticle(
                        id="unknown_auto",
                        title="Alifbo aniqlanmadi (Kirill yoki Lotin deb yozing)",
                        input_message_content=InputTextMessageContent(
                            message_text=f"Alifbo aniqlanmadi: {query_text}"
                        )
                    )
                )
                
    await bot.answer_inline_query(inline_query.id, results, cache_time=5)

# ======================
# Botni ishga tushirish
# ======================

dp.include_router(router)

async def main():
    logger.info("Bot ishga tushdi...")
    if not gemini_client:
        logger.error("!!! GEMINI CLIENT ISHLAMAYAPTI. IMLO TEKSHIRISH FUNKSIYASI O'CHIRILGAN !!!")
    
    # Inline rejimni aktivlashtirish (BotFather orqali ham aktivlashtirish kerak!)
    # BotFather'da /setinline buyrug'ini bosing.
    logger.info("Inline rejim tayyor. @lotinifybot [lotin/kirill] matn ko'rinishida foydalanish mumkin.")

    await dp.start_polling(bot)

if __name__ == "__main__":
    try:
        asyncio.run(main())
    except KeyboardInterrupt:
        logger.info("Bot o'chirildi.")
    except Exception as e:
        logger.exception(f"Kutilmagan xato: {e}")
